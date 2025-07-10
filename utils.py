import pandas as pd
import pptx
import docx
from PyPDF2 import PdfReader

def extract_text_from_file(file):
    if file.name.endswith(".pdf"):
        reader = PdfReader(file)
        return "\n".join([page.extract_text() for page in reader.pages])
    elif file.name.endswith(".pptx"):
        prs = pptx.Presentation(file)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif file.name.endswith(".csv"):
        df = pd.read_csv(file)
        return df.to_string()
    elif file.name.endswith(".docx"):
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file.name.endswith(".txt") or file.name.endswith(".md"):
        return file.read().decode()
    else:
        return ""
